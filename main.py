import streamlit as st
import fitz  # PyMuPDF
from io import BytesIO
from langchain_community.chat_models import ChatOpenAI
from langchain.schema.messages import HumanMessage, AIMessage
import pandas as pd
from docx import Document
import easyocr
from PIL import Image
from pptx import Presentation
import csv

# Initialize LangChain models
chain_gpt_35 = ChatOpenAI(model="gpt-3.5-turbo")
chain_gpt_vision = ChatOpenAI(model="gpt-4-o-mini")

# Initialize EasyOCR reader
ocr_reader = easyocr.Reader(['en'])

# Initialize session state attributes
if 'history' not in st.session_state:
    st.session_state.history = []
if 'files_uploaded' not in st.session_state:
    st.session_state.files_uploaded = []
if 'file_details' not in st.session_state:
    st.session_state.file_details = {}

# Title
st.set_page_config(page_title="RAG File Chatbot", page_icon=":robot:")
st.title("RAG File Chatbot")

# File uploader
files = st.file_uploader("Please Upload file", type=["png", "jpg", "jpeg", "pdf", "txt", "docx", "doc", "ppt", "pptx", "xls", "xlsx", "csv"], accept_multiple_files=True)

# Detect if files have been removed and clear history if so
if files != st.session_state.files_uploaded:
    st.session_state.files_uploaded = files
    st.session_state.history = []  # Clear chat history
    st.session_state.file_details = {}  # Clear file details

# Display uploaded file details
if st.session_state.files_uploaded:
    st.write("### Uploaded Files")
    for i, file in enumerate(st.session_state.files_uploaded):
        file_name = file.name
        st.write(f"file{i+1}: {file_name}")

def process_files(files):
    data_found = False
    for file in files:
        file_type = file.type
        if file_type in ['image/png', 'image/jpeg']:
            data_found |= img_reader(file)
        elif file_type == 'application/pdf':
            data_found |= pdf_reader(file)
        elif file_type == 'text/plain':
            data_found |= txt_reader(file)
        elif file_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
            data_found |= doc_reader(file)
        elif file_type in ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
            data_found |= excel_reader(file)
        elif file_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
            data_found |= ppt_reader(file)
        elif file_type == 'text/csv':
            data_found |= csv_reader(file)
    
    if not data_found:
        st.session_state.history.append({"user": "", "assistant": "The provided files do not meet your requirements."})

def img_reader(file):
    img = Image.open(file)
    text = ocr_reader.readtext(img, detail=0, paragraph=False)
    text = "\n".join(text)
    if text.strip() == "":
        return False
    prompt = f"Answer users query:\n\n{text}\n\n:"
    response = chain_gpt_vision.invoke([HumanMessage(content=prompt)])
    st.session_state.history.append({"user": "Image file", "assistant": response.content})
    return True

def pdf_reader(file):
    try:
        file_stream = BytesIO(file.read())
        if file_stream.getbuffer().nbytes == 0:
            st.session_state.history.append({"user": "", "assistant": "The provided PDF file is empty."})
            return False
        
        pdf_document = fitz.open(stream=file_stream, filetype="pdf")
        full_text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            full_text += page.get_text()
        pdf_document.close()
        
        if full_text.strip() == "":
            st.session_state.history.append({"user": "", "assistant": "The provided PDF file does not contain readable text."})
            return False
        
        prompt = f"Answer users query:\n\n{full_text}\n\n:"
        response = chain_gpt_35.invoke([HumanMessage(content=prompt)])
        st.session_state.history.append({"user": "PDF file", "assistant": response.content})
        return True

    except fitz.EmptyFileError:
        st.session_state.history.append({"user": "", "assistant": "The provided PDF file is empty or corrupted."})
        return False

def txt_reader(file):
    text = file.read().decode("utf-8")
    if text.strip() == "":
        return False
    prompt = f"Answer the following Query for the user:\n\n{text}\n\n:"
    response = chain_gpt_35.invoke([HumanMessage(content=prompt)])
    st.session_state.history.append({"user": "Text file", "assistant": response.content})
    return True

def doc_reader(file):
    doc = Document(BytesIO(file.read()))
    text = "\n".join([para.text for para in doc.paragraphs])
    if text.strip() == "":
        return False
    prompt = f"Answer users query:\n\n{text}\n\n:"
    response = chain_gpt_35.invoke([HumanMessage(content=prompt)])
    st.session_state.history.append({"user": "Word document", "assistant": response.content})
    return True

def excel_reader(file):
    df = pd.read_excel(file)
    text = df.to_string()
    if text.strip() == "":
        return False
    prompt = f"Answer users query:\n\n{text}\n\n:"
    response = chain_gpt_35.invoke([HumanMessage(content=prompt)])
    st.session_state.history.append({"user": "Excel file", "assistant": response.content})
    return True

def ppt_reader(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    if text.strip() == "":
        return False
    prompt = f"Answer users query:\n\n{text}\n\n:"
    response = chain_gpt_35.invoke([HumanMessage(content=prompt)])
    st.session_state.history.append({"user": "PowerPoint file", "assistant": response.content})
    return True

def csv_reader(file):
    df = pd.read_csv(file)
    text = df.to_string()
    if text.strip() == "":
        return False
    prompt = f"Answer users query:\n\n{text}\n\n:"
    response = chain_gpt_35.invoke([HumanMessage(content=prompt)])
    st.session_state.history.append({"user": "CSV file", "assistant": response.content})
    return True

def display_history():
    for message in st.session_state.history:
        st.write(f"**User:** {message['user']}")
        st.write(f"**Assistant:** {message['assistant']}")

# Main chat layout
chat_container = st.container()
input_container = st.container()

with input_container:
    query = st.text_area("Ask your AI Assistant", max_chars=10000, height=200, key="query_input")

    if st.button("Send", use_container_width=True):
        if query:
            # Add user query to history
            st.session_state.history.append({"user": query, "assistant": "Processing..."})

            # Process files if any are uploaded
            if st.session_state.files_uploaded:
                process_files(st.session_state.files_uploaded)
            
            # Get the assistant's response
            history_messages = [HumanMessage(content=msg['user']) for msg in st.session_state.history]
            history_messages.append(HumanMessage(content=query))
            
            response = chain_gpt_35.invoke(history_messages)
            st.session_state.history[-1]["assistant"] = response.content  # Update latest response

    with chat_container:
        st.write("### Conversation History")
        display_history()
